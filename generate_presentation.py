#!/usr/bin/env python3
"""
Patrick Kilzi AI - Générateur de Présentations
Usage: python3 generate_presentation.py
"""

import os
import sys
import json
import anthropic
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import datetime


# ─── PALETTE INSIGHTER ───────────────────────────────────────────────────────
NAVY       = RGBColor(0x1A, 0x3C, 0x5E)
TEAL       = RGBColor(0x0D, 0x94, 0x88)
AMBER      = RGBColor(0xF5, 0x9E, 0x0B)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG   = RGBColor(0xF8, 0xFA, 0xFC)
DARK_TEXT  = RGBColor(0x1E, 0x29, 0x3B)
MUTED_TEXT = RGBColor(0x64, 0x74, 0x8B)
CARD_BG    = RGBColor(0xFF, 0xFF, 0xFF)
BORDER     = RGBColor(0xE2, 0xE8, 0xF0)


def hex_rgb(r, g, b):
    return RGBColor(r, g, b)


def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=None):
    from pptx.util import Pt as PtUtil
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h, font_size=14, bold=False, color=None,
             align=PP_ALIGN.LEFT, italic=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    if color:
        run.font.color.rgb = color
    return txBox


def slide_title_page(prs, title, subtitle, client, date_str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Dark navy background
    add_rect(slide, 0, 0, 10, 5.625, NAVY)

    # Teal accent bar left
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)

    # Amber accent bottom bar
    add_rect(slide, 0, 4.8, 10, 0.12, AMBER)

    # Company name top right
    add_text(slide, "INSIGHTER", 7.5, 0.2, 2.3, 0.4,
             font_size=10, bold=True, color=TEAL, align=PP_ALIGN.RIGHT, font_name="Calibri")

    # Main title
    add_text(slide, title.upper(), 0.5, 1.5, 9, 1.0,
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT, font_name="Calibri")

    # Separator line
    add_rect(slide, 0.5, 2.65, 1.2, 0.04, TEAL)

    # Subtitle
    add_text(slide, subtitle, 0.5, 2.8, 9, 0.6,
             font_size=16, bold=False, color=RGBColor(0xCA, 0xDC, 0xFC),
             align=PP_ALIGN.LEFT, font_name="Calibri")

    # Client + date bottom
    add_text(slide, f"Client : {client}   |   {date_str}", 0.5, 4.3, 9, 0.35,
             font_size=11, color=MUTED_TEXT, align=PP_ALIGN.LEFT, font_name="Calibri")


def slide_section_divider(prs, section_title, section_num):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, LIGHT_BG)
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)
    add_rect(slide, 0, 5.3, 10, 0.08, BORDER)

    # Number
    add_text(slide, str(section_num).zfill(2), 0.5, 1.6, 2, 1.2,
             font_size=72, bold=True, color=RGBColor(0xE2, 0xE8, 0xF0),
             align=PP_ALIGN.LEFT, font_name="Calibri")

    # Title
    add_text(slide, section_title.upper(), 0.5, 2.2, 9, 0.8,
             font_size=28, bold=True, color=NAVY, align=PP_ALIGN.LEFT, font_name="Calibri")

    add_text(slide, "INSIGHTER", 7.5, 5.2, 2.3, 0.3,
             font_size=9, bold=True, color=MUTED_TEXT, align=PP_ALIGN.RIGHT, font_name="Calibri")


def slide_executive_summary(prs, summary_points):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, LIGHT_BG)
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)
    add_rect(slide, 0, 5.3, 10, 0.08, BORDER)

    add_text(slide, "RÉSUMÉ EXÉCUTIF", 0.35, 0.18, 9, 0.45,
             font_size=11, bold=True, color=TEAL, font_name="Calibri")
    add_rect(slide, 0.35, 0.7, 9.3, 0.04, BORDER)

    y = 0.85
    for i, point in enumerate(summary_points[:4]):
        # Card background
        add_rect(slide, 0.35, y, 9.3, 0.85, CARD_BG, BORDER, 0.5)
        # Teal left accent
        add_rect(slide, 0.35, y, 0.07, 0.85, TEAL)
        # Text
        add_text(slide, point, 0.6, y + 0.08, 8.8, 0.7,
                 font_size=13, color=DARK_TEXT, font_name="Calibri")
        y += 0.97

    add_text(slide, "INSIGHTER", 7.5, 5.2, 2.3, 0.3,
             font_size=9, bold=True, color=MUTED_TEXT, align=PP_ALIGN.RIGHT, font_name="Calibri")


def slide_kpis(prs, slide_title, kpis):
    """KPI cards — up to 4 stats"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, LIGHT_BG)
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)
    add_rect(slide, 0, 5.3, 10, 0.08, BORDER)

    add_text(slide, slide_title.upper(), 0.35, 0.18, 9, 0.45,
             font_size=11, bold=True, color=TEAL, font_name="Calibri")
    add_rect(slide, 0.35, 0.7, 9.3, 0.04, BORDER)

    n = min(len(kpis), 4)
    card_w = 9.3 / n - 0.1
    x = 0.35
    for kpi in kpis[:n]:
        add_rect(slide, x, 0.85, card_w, 2.1, CARD_BG, BORDER, 0.5)
        add_rect(slide, x, 0.85, card_w, 0.08, TEAL)

        add_text(slide, kpi.get("value", ""), x + 0.15, 1.05, card_w - 0.3, 0.9,
                 font_size=36, bold=True, color=NAVY, align=PP_ALIGN.CENTER, font_name="Calibri")
        add_text(slide, kpi.get("label", ""), x + 0.1, 1.95, card_w - 0.2, 0.6,
                 font_size=11, color=MUTED_TEXT, align=PP_ALIGN.CENTER, font_name="Calibri")
        if kpi.get("trend"):
            trend_color = TEAL if "+" in kpi["trend"] else RGBColor(0xEF, 0x44, 0x44)
            add_text(slide, kpi["trend"], x + 0.1, 2.55, card_w - 0.2, 0.35,
                     font_size=11, bold=True, color=trend_color,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
        x += card_w + 0.1

    add_text(slide, "INSIGHTER", 7.5, 5.2, 2.3, 0.3,
             font_size=9, bold=True, color=MUTED_TEXT, align=PP_ALIGN.RIGHT, font_name="Calibri")


def slide_content(prs, slide_title, body_text, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, LIGHT_BG)
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)
    add_rect(slide, 0, 5.3, 10, 0.08, BORDER)

    add_text(slide, slide_title.upper(), 0.35, 0.18, 9, 0.45,
             font_size=11, bold=True, color=TEAL, font_name="Calibri")
    add_rect(slide, 0.35, 0.7, 9.3, 0.04, BORDER)

    if subtitle:
        add_text(slide, subtitle, 0.35, 0.78, 9.3, 0.4,
                 font_size=18, bold=True, color=NAVY, font_name="Calibri")
        content_y = 1.3
    else:
        content_y = 0.85

    # Content card
    add_rect(slide, 0.35, content_y, 9.3, 4.1 - (content_y - 0.85), CARD_BG, BORDER, 0.5)
    add_text(slide, body_text, 0.55, content_y + 0.15,
             8.9, 3.6 - (content_y - 0.85),
             font_size=13, color=DARK_TEXT, font_name="Calibri")

    add_text(slide, "INSIGHTER", 7.5, 5.2, 2.3, 0.3,
             font_size=9, bold=True, color=MUTED_TEXT, align=PP_ALIGN.RIGHT, font_name="Calibri")


def slide_two_columns(prs, slide_title, left_title, left_items, right_title, right_items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, LIGHT_BG)
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)
    add_rect(slide, 0, 5.3, 10, 0.08, BORDER)

    add_text(slide, slide_title.upper(), 0.35, 0.18, 9, 0.45,
             font_size=11, bold=True, color=TEAL, font_name="Calibri")
    add_rect(slide, 0.35, 0.7, 9.3, 0.04, BORDER)

    # Left column
    add_rect(slide, 0.35, 0.85, 4.5, 4.1, CARD_BG, BORDER, 0.5)
    add_rect(slide, 0.35, 0.85, 4.5, 0.45, NAVY)
    add_text(slide, left_title, 0.5, 0.9, 4.2, 0.35,
             font_size=12, bold=True, color=WHITE, font_name="Calibri")

    y = 1.45
    for item in left_items[:5]:
        add_rect(slide, 0.45, y, 0.12, 0.12, TEAL)
        add_text(slide, item, 0.65, y - 0.03, 4.0, 0.3,
                 font_size=12, color=DARK_TEXT, font_name="Calibri")
        y += 0.55

    # Right column
    add_rect(slide, 5.15, 0.85, 4.5, 4.1, CARD_BG, BORDER, 0.5)
    add_rect(slide, 5.15, 0.85, 4.5, 0.45, TEAL)
    add_text(slide, right_title, 5.3, 0.9, 4.2, 0.35,
             font_size=12, bold=True, color=WHITE, font_name="Calibri")

    y = 1.45
    for item in right_items[:5]:
        add_rect(slide, 5.25, y, 0.12, 0.12, NAVY)
        add_text(slide, item, 5.45, y - 0.03, 4.0, 0.3,
                 font_size=12, color=DARK_TEXT, font_name="Calibri")
        y += 0.55

    add_text(slide, "INSIGHTER", 7.5, 5.2, 2.3, 0.3,
             font_size=9, bold=True, color=MUTED_TEXT, align=PP_ALIGN.RIGHT, font_name="Calibri")


def slide_recommendations(prs, recommendations):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, LIGHT_BG)
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)
    add_rect(slide, 0, 5.3, 10, 0.08, BORDER)

    add_text(slide, "RECOMMANDATIONS", 0.35, 0.18, 9, 0.45,
             font_size=11, bold=True, color=TEAL, font_name="Calibri")
    add_rect(slide, 0.35, 0.7, 9.3, 0.04, BORDER)

    y = 0.85
    colors = [NAVY, TEAL, AMBER, RGBColor(0x7C, 0x3A, 0xED)]
    for i, reco in enumerate(recommendations[:4]):
        c = colors[i % len(colors)]
        add_rect(slide, 0.35, y, 9.3, 0.88, CARD_BG, BORDER, 0.5)
        add_rect(slide, 0.35, y, 0.5, 0.88, c)
        num_box = slide.shapes.add_textbox(Inches(0.35), Inches(y), Inches(0.5), Inches(0.88))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(i + 1)
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = "Calibri"

        add_text(slide, reco, 1.0, y + 0.1, 8.5, 0.7,
                 font_size=13, color=DARK_TEXT, font_name="Calibri")
        y += 1.0

    add_text(slide, "INSIGHTER", 7.5, 5.2, 2.3, 0.3,
             font_size=9, bold=True, color=MUTED_TEXT, align=PP_ALIGN.RIGHT, font_name="Calibri")


def slide_next_steps(prs, steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, LIGHT_BG)
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)
    add_rect(slide, 0, 5.3, 10, 0.08, BORDER)

    add_text(slide, "PROCHAINES ÉTAPES", 0.35, 0.18, 9, 0.45,
             font_size=11, bold=True, color=TEAL, font_name="Calibri")
    add_rect(slide, 0.35, 0.7, 9.3, 0.04, BORDER)

    n = min(len(steps), 4)
    col_w = 9.3 / n - 0.1
    x = 0.35
    for i, step in enumerate(steps[:n]):
        add_rect(slide, x, 0.85, col_w, 3.6, CARD_BG, BORDER, 0.5)
        add_rect(slide, x, 0.85, col_w, 0.08, [TEAL, NAVY, AMBER, RGBColor(0x7C, 0x3A, 0xED)][i])

        # Step number
        add_text(slide, str(i + 1), x + col_w / 2 - 0.3, 1.1, 0.6, 0.6,
                 font_size=28, bold=True, color=RGBColor(0xE2, 0xE8, 0xF0),
                 align=PP_ALIGN.CENTER, font_name="Calibri")

        # Step title
        add_text(slide, step.get("title", ""), x + 0.1, 1.8, col_w - 0.2, 0.5,
                 font_size=12, bold=True, color=NAVY,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

        # Step description
        add_text(slide, step.get("description", ""), x + 0.1, 2.4, col_w - 0.2, 1.8,
                 font_size=11, color=MUTED_TEXT,
                 align=PP_ALIGN.CENTER, font_name="Calibri")

        if step.get("deadline"):
            add_text(slide, step["deadline"], x + 0.1, 4.1, col_w - 0.2, 0.3,
                     font_size=10, bold=True, color=TEAL,
                     align=PP_ALIGN.CENTER, font_name="Calibri")
        x += col_w + 0.1

    add_text(slide, "INSIGHTER", 7.5, 5.2, 2.3, 0.3,
             font_size=9, bold=True, color=MUTED_TEXT, align=PP_ALIGN.RIGHT, font_name="Calibri")


def slide_closing(prs, thank_you_msg, contact):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, 10, 5.625, NAVY)
    add_rect(slide, 0, 0, 0.12, 5.625, TEAL)
    add_rect(slide, 0, 4.8, 10, 0.12, AMBER)

    add_text(slide, "INSIGHTER", 7.5, 0.2, 2.3, 0.4,
             font_size=10, bold=True, color=TEAL, align=PP_ALIGN.RIGHT, font_name="Calibri")

    add_text(slide, thank_you_msg, 0.5, 1.8, 9, 0.9,
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER, font_name="Calibri")

    add_rect(slide, 4.0, 2.85, 2.0, 0.06, TEAL)

    add_text(slide, contact, 0.5, 3.1, 9, 0.5,
             font_size=13, color=RGBColor(0xCA, 0xDC, 0xFC),
             align=PP_ALIGN.CENTER, font_name="Calibri")


# ─── CLAUDE API ──────────────────────────────────────────────────────────────

def generate_slide_structure(project_info: str, client: str) -> dict:
    """Appelle Claude pour structurer le contenu en slides."""
    ai = anthropic.Anthropic()

    prompt = f"""Tu es un expert en Business Intelligence et présentation executive.
Génère la structure complète d'une présentation professionnelle en JSON strict.

Projet / Informations :
{project_info}

Client : {client}

Retourne UNIQUEMENT un JSON valide avec cette structure exacte :
{{
  "title": "Titre principal de la présentation",
  "subtitle": "Sous-titre ou description courte",
  "executive_summary": [
    "Point clé 1 (une phrase percutante)",
    "Point clé 2",
    "Point clé 3",
    "Point clé 4"
  ],
  "kpis": [
    {{"value": "€2.4M", "label": "Chiffre d'affaires", "trend": "+12%"}},
    {{"value": "87%", "label": "Taux de satisfaction", "trend": "+5%"}},
    {{"value": "3", "label": "Nouveaux marchés", "trend": null}}
  ],
  "sections": [
    {{
      "title": "Titre de la section",
      "subtitle": "Sous-titre optionnel",
      "content": "Contenu détaillé de cette section (3-5 phrases).",
      "type": "content"
    }},
    {{
      "title": "Analyse comparative",
      "left_title": "Points forts",
      "left_items": ["Item 1", "Item 2", "Item 3"],
      "right_title": "Points d'amélioration",
      "right_items": ["Item 1", "Item 2", "Item 3"],
      "type": "two_columns"
    }}
  ],
  "recommendations": [
    "Recommandation 1 concrète et actionnable",
    "Recommandation 2",
    "Recommandation 3",
    "Recommandation 4"
  ],
  "next_steps": [
    {{"title": "Étape 1", "description": "Description courte", "deadline": "Semaine 1"}},
    {{"title": "Étape 2", "description": "Description courte", "deadline": "Semaine 2"}},
    {{"title": "Étape 3", "description": "Description courte", "deadline": "Semaine 3"}},
    {{"title": "Étape 4", "description": "Description courte", "deadline": "Semaine 4"}}
  ],
  "closing_message": "Message de clôture (ex: Merci pour votre confiance)",
  "contact": "contact@insighter.com | +33 6 XX XX XX XX"
}}

Important :
- Le contenu doit être professionnel et adapté à un dirigeant
- Utilise les informations fournies pour personnaliser le contenu
- Minimum 2 sections, maximum 4 sections
- Les KPIs doivent être des vrais chiffres si fournis, sinon invente des exemples cohérents
- Réponds UNIQUEMENT avec le JSON, sans texte avant ou après
"""

    message = ai.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=4000,
        messages=[{"role": "user", "content": prompt}]
    )

    raw = message.content[0].text.strip()
    # Remove markdown code blocks if present
    if raw.startswith("```"):
        raw = raw.split("```")[1]
        if raw.startswith("json"):
            raw = raw[4:]
        raw = raw.strip()

    return json.loads(raw)


# ─── BUILD PRESENTATION ──────────────────────────────────────────────────────

def build_presentation(data: dict, client: str, output_path: str):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)

    date_str = datetime.date.today().strftime("%d %B %Y")

    # 1. Slide titre
    slide_title_page(prs, data["title"], data["subtitle"], client, date_str)

    # 2. Résumé exécutif
    slide_executive_summary(prs, data["executive_summary"])

    # 3. KPIs
    if data.get("kpis"):
        slide_kpis(prs, "Indicateurs clés de performance", data["kpis"])

    # 4. Sections de contenu
    for i, section in enumerate(data.get("sections", []), 1):
        if i == 1:
            slide_section_divider(prs, section["title"], i)

        if section.get("type") == "two_columns":
            slide_two_columns(
                prs, section["title"],
                section.get("left_title", "Colonne 1"),
                section.get("left_items", []),
                section.get("right_title", "Colonne 2"),
                section.get("right_items", [])
            )
        else:
            slide_content(prs, section["title"],
                          section.get("content", ""),
                          section.get("subtitle"))

    # 5. Recommandations
    if data.get("recommendations"):
        slide_section_divider(prs, "Recommandations", len(data.get("sections", [])) + 1)
        slide_recommendations(prs, data["recommendations"])

    # 6. Prochaines étapes
    if data.get("next_steps"):
        slide_next_steps(prs, data["next_steps"])

    # 7. Clôture
    slide_closing(prs, data.get("closing_message", "Merci pour votre confiance"),
                  data.get("contact", "contact@insighter.com"))

    prs.save(output_path)


# ─── MAIN ────────────────────────────────────────────────────────────────────

def main():
    print("\n" + "="*60)
    print("  INSIGHTER AI — Générateur de Présentations")
    print("="*60 + "\n")

    # Vérification clé API
    if not os.environ.get("ANTHROPIC_API_KEY"):
        print("❌ Erreur : ANTHROPIC_API_KEY non définie.")
        print("   Exécute : export ANTHROPIC_API_KEY='ta-clé-ici'\n")
        sys.exit(1)

    # Input client
    client = input("📋 Nom du client : ").strip()
    if not client:
        client = "Client"

    print("\n📝 Décris le projet/analyse (appuie sur Entrée 2x pour terminer) :")
    lines = []
    while True:
        line = input()
        if line == "" and lines and lines[-1] == "":
            break
        lines.append(line)
    project_info = "\n".join(lines).strip()

    if not project_info:
        print("❌ Description vide. Arrêt.")
        sys.exit(1)

    # Nom du fichier de sortie
    safe_client = client.replace(" ", "_").replace("/", "-")
    timestamp = datetime.datetime.now().strftime("%Y%m%d_%H%M")
    output_path = f"Presentation_{safe_client}_{timestamp}.pptx"

    print("\n⏳ Génération du contenu par l'IA...")
    try:
        data = generate_slide_structure(project_info, client)
        print(f"✅ Structure générée : {data['title']}")
    except Exception as e:
        print(f"❌ Erreur API Claude : {e}")
        sys.exit(1)

    print("⏳ Création des slides...")
    try:
        build_presentation(data, client, output_path)
        print(f"\n✅ Présentation créée : {output_path}")
        print(f"   Slides : {len(data.get('sections', [])) + 6} slides")
        print(f"   Client : {client}")
        print(f"   Date   : {datetime.date.today().strftime('%d/%m/%Y')}\n")
    except Exception as e:
        print(f"❌ Erreur création PPTX : {e}")
        import traceback
        traceback.print_exc()
        sys.exit(1)


if __name__ == "__main__":
    main()
